"""シルバー層: テキスト抽出・クレンジング・PII・重複排除(純関数中心)。

- extract_text / clean_text / redact_pii / make_document_id は spark 非依存の純関数
  → ローカルでユニットテスト可能。
- add_extracted_text / dedup_documents は DataFrame を受け取る Spark ヘルパー。
- パーサー(pypdf / python-pptx)は関数内で遅延import。
"""

import hashlib
import re

from pyspark.sql import DataFrame
from pyspark.sql.functions import col, udf
from pyspark.sql.types import StringType

from lib.unstructured_rag.config import CORRUPT_COLUMN, detect_source_type

# ---- PII 検出用の簡易パターン(デモ用。実運用はより厳密に) ----
_PII_PATTERNS = {
    "email": re.compile(r"[\w.+-]+@[\w-]+\.[\w.-]+"),
    "phone": re.compile(r"\b0\d{1,4}-?\d{1,4}-?\d{3,4}\b"),
}


# ============================================================
# 純関数(テスト対象)
# ============================================================
def make_document_id(file_path: str) -> str:
    """file_path から決定的な document_id を生成する。

    決定的にすることで、再取り込み時も同じIDになり MERGE(Upsert)しやすい。
    """
    return hashlib.sha256(file_path.encode("utf-8")).hexdigest()[:32]


def extract_text(content: bytes, source_type: str) -> str:
    """生バイナリから source_type に応じてテキストを抽出する(純関数)。

    未対応 source_type や抽出失敗時は例外を送出する(呼び出し側で _corrupt 退避)。
    """
    if source_type == "pdf":
        return _extract_pdf(content)
    if source_type == "pptx":
        return _extract_pptx(content)
    if source_type == "md":
        return content.decode("utf-8", errors="replace")
    raise ValueError(f"unsupported source_type: {source_type}")


def _extract_pdf(content: bytes) -> str:
    import io

    from pypdf import PdfReader  # 遅延import

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_pptx(content: bytes) -> str:
    import io

    from pptx import Presentation  # 遅延import (python-pptx)

    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n\n".join(t for t in texts if t)


def clean_text(text: str) -> str:
    """余分な空白・制御文字を正規化する(純関数)。"""
    if text is None:
        return ""
    text = text.replace("\x00", "")
    # 連続する空行を1つに、行末の空白を除去
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def redact_pii(text: str) -> str:
    """メール・電話番号などのPIIをマスクする(純関数)。"""
    if not text:
        return text
    for label, pattern in _PII_PATTERNS.items():
        text = pattern.sub(f"[REDACTED_{label.upper()}]", text)
    return text


# ============================================================
# Spark ヘルパー(DataFrame を受け取る)
# ============================================================
def _extract_and_clean(content: bytes, file_path: str) -> str:
    """UDF本体: 抽出→クレンジング→PIIマスク。失敗時は None(→_corrupt退避)。"""
    source_type = detect_source_type(file_path)
    if source_type is None or content is None:
        return None
    try:
        raw = extract_text(bytes(content), source_type)
        return redact_pii(clean_text(raw))
    except Exception:
        return None


def add_extracted_text(df: DataFrame) -> DataFrame:
    """Bronze(file_path, binary_content) から Silver列を生成する。

    抽出失敗行は extracted_text=null / _corrupt にエラー印を付ける。
    """
    extract_udf = udf(_extract_and_clean, StringType())
    source_type_udf = udf(detect_source_type, StringType())

    return (
        df.withColumn("source_type", source_type_udf(col("file_path")))
        .withColumn("extracted_text", extract_udf(col("binary_content"), col("file_path")))
        .withColumn(
            CORRUPT_COLUMN,
            col(CORRUPT_COLUMN)  # Bronze由来のnullを保持しつつ、抽出失敗を印字
            if CORRUPT_COLUMN in df.columns
            else col("extracted_text"),
        )
    )


def dedup_documents(df: DataFrame, key_col: str = "extracted_text") -> DataFrame:
    """ほぼ同一ドキュメントを除外する(抽出テキストのハッシュで重複排除)。"""
    from pyspark.sql.functions import sha2

    return (
        df.withColumn("_text_hash", sha2(col(key_col), 256))
        .dropDuplicates(["_text_hash"])
        .drop("_text_hash")
    )
